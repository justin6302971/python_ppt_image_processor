from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE
from pptx.util import Pt, Cm, Inches
from pptx.dml.color import ColorFormat, RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE


from os import listdir, makedirs, getenv
from os.path import isfile, join, exists, dirname, splitext
from datetime import datetime as d
import constants
import sys
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from PIL import Image
from pillow_heif import register_heif_opener

register_heif_opener()


# dotenv setup
from dotenv import load_dotenv




# determine if application is a script file or frozen exe
if getattr(sys, 'frozen', False):
    application_path = dirname(sys.executable)
elif __file__:
    application_path = dirname(__file__)

dotenv_path = join(application_path, '.env')
load_dotenv(dotenv_path)

is_insert_checking_icon = getenv("IS_INSERT_CHECKING_ICON") == 'True'


font_size = int(getenv("FONT_SIZE"))
if font_size == "" or font_size is None:
    font_size = 12

font_type = getenv("FONT_TYPE")
if font_type == "" or font_type is None:
    font_type = "Microsoft JhengHei"

font_alignment_center = getenv("FONT_ALIGNMENT_CENTER") == 'True'
font_alignment = PP_ALIGN.LEFT
if font_alignment_center:
    font_alignment = PP_ALIGN.CENTER


output_file_name = getenv("FILE_NAME")
if output_file_name == "" or output_file_name is None:
    output_file_name = constants.FILE_NAME

template_file_name = getenv(
    "TEMPLATE_FILE_NAME") or constants.TEMPLATE_FILE_NAME
if template_file_name == "" or template_file_name is None:
    template_file_name = constants.TEMPLATE_FILE_NAME

custom_image_dir_name = getenv("CUSTOM_IMAGE_DIR")
if custom_image_dir_name == "" or custom_image_dir_name is None:
    custom_image_dir_name = constants.IMAGE_DIR
    image_dir = join(application_path, constants.IMAGE_DIR)
else:
    image_dir = join(application_path, constants.IMAGE_DIR,
                     custom_image_dir_name)


date = d.now()
output_file_name_with_date = f'{date.strftime("%Y%m%d%H%M%S")}_{output_file_name}'

print(f'assigned output file name:{output_file_name_with_date}')


# check if template exists
template_dir = join(application_path, constants.TEMPLATE_DIR)

is_template_dir_exist = exists(template_dir)

if not is_template_dir_exist:
    makedirs(template_dir)

template_file_path = join(template_dir, template_file_name)

print(f'assigned template name:{template_file_name}')

is_template_file_exist = exists(template_file_path)
if not is_template_file_exist:
    raise Exception(
        f'template file "{template_file_name}" doesn\'t exist in the follow directory: {template_dir}')



is_image_dir_exist = exists(image_dir)

if not is_image_dir_exist:
    makedirs(image_dir)

processed_image_dir = join(image_dir, constants.PROCESS_IMAGE_DIR)

is_processed_image_dir_exist = exists(processed_image_dir)

if not is_processed_image_dir_exist:
    makedirs(processed_image_dir)


image_files = [f for f in listdir(image_dir) if isfile(
    join(image_dir, f)) and not f.startswith('.')]

if len(image_files) == 0:
    raise Exception(
        f'no images exist in the follow directory: {image_dir}')

print("start ------- inserting the images to ppt in ascending order")


sorted_image_files = sorted(image_files)

paginated_sorted_image_files = [sorted_image_files[i:i+4]
                                for i in range(0, len(sorted_image_files), 4)]


prs = Presentation(template_file_path)

# Adding intro slides
current_datetime = date.strftime("%Y-%m-%d %H:%M:%S")
note_slide = prs.slides
for slide in note_slide:
    slide.shapes.title.text = "generated image slides"
    slide.placeholders[1].text = current_datetime

custom_slide_layout = prs.slide_layouts[11]


for item_array in paginated_sorted_image_files:
    slide = prs.slides.add_slide(custom_slide_layout)
    try:
        images_placeholder_idx_arr = []
        text_placeholder_idx_arr = []

        for shape in slide.placeholders:
            if shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
                images_placeholder_idx_arr.append(
                    {"idx": shape.placeholder_format.idx, "name": shape.name})
            if shape.placeholder_format.type == PP_PLACEHOLDER.OBJECT:
                text_placeholder_idx_arr.append(
                    {"idx": shape.placeholder_format.idx, "name": shape.name})
            # print('%d %s %s' % (shape.placeholder_format.idx,
            #       shape.name, shape.placeholder_format.type))

        images_placeholder_idx_arr_sorted_by_name = sorted(
            images_placeholder_idx_arr, key=lambda x: x['name'])
        text_placeholder_idx_arr_sorted_by_name = sorted(
            text_placeholder_idx_arr, key=lambda x: x['name'])

        for index in range(0, len(item_array)):
            # print(idx)
            image_name = item_array[index]
            image_path = join(image_dir, image_name)

            im = Image.open(image_path)
            file_name_without_ext= splitext(image_name)[0]
            processed_file_name=file_name_without_ext+".JPEG"
            processed_image_path = join(processed_image_dir, processed_file_name )
            
            im.save(processed_image_path)
            

            placeholder = slide.placeholders[images_placeholder_idx_arr_sorted_by_name[index]['idx']]
            picture = placeholder.insert_picture(processed_image_path)

            placeholder_text = slide.placeholders[text_placeholder_idx_arr_sorted_by_name[index]['idx']]
            processed_file_name = splitext(image_name)[0].replace("=", ":")

            tf = placeholder_text.text_frame

            p = tf.paragraphs[0]
            p.alignment = font_alignment
            p.text = processed_file_name
            p.font.size = Pt(font_size)
            p.font.name = font_type

        if is_insert_checking_icon:

            circle_width = Cm(1.3)
            circle_height = Cm(1.3)

            top_1 = Cm(5.87)
            left_1 = Cm(1.77)

            circle_shape_1 = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, left_1, top_1, circle_width, circle_height)
            circle_shape_1.fill.solid()
            circle_shape_1.fill.patterned()
            circle_shape_1.fill.fore_color.rgb = RGBColor(255, 255, 255)
            circle_shape_1.fill.back_color.rgb = RGBColor(255, 255, 255)
            circle_shape_1.line.color.rgb = RGBColor(255, 255, 255)

            top_2 = Cm(5.87)
            left_2 = Cm(10.91)

            circle_shape_2 = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, left_2, top_2, circle_width, circle_height)
            circle_shape_2.fill.solid()
            circle_shape_2.fill.patterned()
            circle_shape_2.fill.fore_color.rgb = RGBColor(255, 255, 255)
            circle_shape_2.fill.back_color.rgb = RGBColor(255, 255, 255)
            circle_shape_2.line.color.rgb = RGBColor(255, 255, 255)

            top_3 = Cm(18.55)
            left_3 = Cm(1.77)

            circle_shape_3 = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, left_3, top_3, circle_width, circle_height)
            circle_shape_3.fill.solid()
            circle_shape_3.fill.patterned()
            circle_shape_3.fill.fore_color.rgb = RGBColor(255, 255, 255)
            circle_shape_3.fill.back_color.rgb = RGBColor(255, 255, 255)
            circle_shape_3.line.color.rgb = RGBColor(255, 255, 255)

            top_4 = Cm(18.55)
            left_4 = Cm(10.91)

            circle_shape_4 = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, left_4, top_4, circle_width, circle_height)
            circle_shape_4.fill.solid()
            circle_shape_4.fill.patterned()
            circle_shape_4.fill.fore_color.rgb = RGBColor(255, 255, 255)
            circle_shape_4.fill.back_color.rgb = RGBColor(255, 255, 255)
            circle_shape_4.line.color.rgb = RGBColor(255, 255, 255)

    except Exception as e:
        print("image insertion issues")

print("end ------- inserting the images to ppt in ascending order")

print("start ------- inserting the pagenumber to ppt slides")


SH = prs.slide_height
SW = prs.slide_width
OutsideMargin_x = Cm(1)
OutsideMargin_y = Cm(0.5)

# i = 0
for slide in prs.slides:
    # i = i + 1
    txBox = slide.shapes.add_textbox(0, 0, 2, 3)
    txBox.height = Cm(0.86)
    txBox.width = Cm(1.63)
    txBox.top = SH - OutsideMargin_y - txBox.height
    txBox.left = SW - OutsideMargin_x - txBox.width
    tf = txBox.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Arial'
    p.font.color.rgb =RGBColor(255, 255, 255)
    p.font.size = Pt(25.9)
    p.font.bold=True
    # run = p.add_run()
    # run.text = str(i)
    # run.font.size = Pt(25.9)

    # ---get a textbox paragraph element---
    pe =p._p

    # ---add fld element---
    fld_xml = (
        '<a:fld %s id="{1F4E2DE4-8ADA-4D4E-9951-90A1D26586E7}" type="slidenum">\n'
        '  <a:rPr lang="en-US" smtClean="0"/>\n'
        '  <a:t>2</a:t>\n'
        '</a:fld>\n' % nsdecls("a")
    )
    fld = parse_xml(fld_xml)
    pe.append(fld)

print("end ------- inserting the pagenumber to ppt slides")

output_dir = join(application_path, constants.OUTPUT_DIR)

is_output_dir_exist = exists(output_dir)

if not is_output_dir_exist:
    makedirs(output_dir)

print(
    f'generate file"{output_file_name_with_date}" to this directory: {output_dir}')

output_path = join(output_dir, output_file_name_with_date)

prs.save(output_path)